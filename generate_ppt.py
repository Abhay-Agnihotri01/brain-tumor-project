from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.name = 'Calibri'
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(41, 128, 185) # Blue title

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content
    
    # Format paragraph bullet points
    for p in tf.paragraphs:
        p.font.size = Pt(20)
        p.font.name = 'Calibri'
        p.space_after = Pt(10)

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Lightweight GAN-VAE Hybrid Model\nfor Real-Time Brain Tumor Detection on Edge Devices"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(41, 128, 185)
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text = "Comprehensive Research & Evaluation\nPresenter Name"

    # Slide 2: Introduction
    add_slide(prs, "Introduction to Research",
              "• Brain tumors are deadly; early, precise diagnosis via MRI is crucial to improve prognosis.\n"
              "• While CNNs are highly effective, they often demand substantial computational resources.\n"
              "• This limits deployment in real-time and resource-constrained environments (e.g., small clinics, edge devices).\n"
              "• Furthermore, genetic characteristics influence tumor susceptibility, requiring multidimensional, fast analysis.\n"
              "• Objective: Introduce an efficient VAE-GAN hybrid framework optimized for low-power hardware like the Raspberry Pi.")

    # Slide 3: Related Work & Motivation
    add_slide(prs, "Related Work & Motivation",
              "• Evolution: Basic image processing \u2192 CNNs \u2192 Generative Models (GANs/VAEs).\n"
              "• Prior generative models (e.g., GANs for synthetic data) are highly accurate but often too resource-heavy.\n"
              "• The Gap: Most state-of-the-art models (e.g., GAN-Transformer) require powerful GPUs.\n"
              "• Solution: A model that is small, fast, and deployable on edge devices, maintaining high accuracy through generative feature fusion.")

    # Slide 4: Proposed Methodology - Architecture
    add_slide(prs, "Proposed Methodology: System Architecture",
              "• Hybrid Framework integrates 4 main components for efficiency and accuracy:\n"
              "  1. VAE Encoder: Compresses MRI images to extract stable latent features.\n"
              "  2. GAN Generator: Synthesizes high-fidelity tumor images from latent vectors.\n"
              "  3. Discriminator: Enforces realistic image generation via adversarial training.\n"
              "  4. Lightweight Classifier: A MobileNetV2-based CNN predicting tumors efficiently.\n\n"
              "[PLACE IMAGE HERE: Figure 1. Proposed lightweight GAN-VAE hybrid architecture diagram]")

    # Slide 5: Experimental Setup
    add_slide(prs, "Experimental Setup & Edge Simulation",
              "• Dataset: Synthetic MRI data (BraTS emulation), normalized and resized to 64x64 pixels.\n"
              "• Data Augmentation: Applied rotation, flipping, and Gaussian noise to prevent overfitting.\n"
              "• Base Environment: Simulated edge environment on a quad-core ARM Cortex-A72 (comparable to Raspberry Pi 4), CPU-only.\n"
              "• Optimization: Models were trained conventionally on GPU, then quantized and deployed using ONNX Runtime/TensorFlow Lite for edge.")

    # -------------------------------------------------------------
    # STANDALONE VAE SECTION
    # -------------------------------------------------------------
    add_slide(prs, "Model 1: Variational Autoencoder (VAE) Overview",
              "• Architecture: Deep Convolutional Encoder & Decoder.\n"
              "• Function: Learns a continuous latent space (mu, log_var) and reconstructs input images.\n"
              "• Classifier: A 3-block CNN (32->64->128 filters) with heavy 0.4 Dropout.\n"
              "• Training Loss = Reconstruction Loss (MSE) + KL-Divergence + Classification Loss (BCE).\n"
              "• Purpose: Evaluates raw feature-extraction capabilities without adversarial loss.")

    add_slide(prs, "VAE Results: Training & Testing Loss",
              "[PLACE IMAGE HERE: VAE Model - Training vs Testing Loss Curve]\n\n"
              "Notes: The dual-line graph shows the VAE reconstruction/KL loss convergence over 25 epochs.")

    add_slide(prs, "VAE Results: Confusion Matrix & ROC AUC",
              "[PLACE IMAGE HERE: VAE Model - Confusion Matrix Heatmap]\n\n"
              "[PLACE IMAGE HERE: VAE Model - ROC Curve (AUC)]\n\n"
              "Notes: The Seaborn heatmap shows False/True Positives. The ROC AUC curve displays diagnostic capability.")

    add_slide(prs, "VAE Results: Efficiency & Final Metrics",
              "[PLACE IMAGE HERE: VAE Model - Accuracy vs Model Size Line Graph]\n\n"
              "[PLACE IMAGE HERE: VAE Model - Final Metrics Snippet (Accuracy, Dice, Size, Inference)]\n\n"
              "Notes: Simulated ablation line graph comparing size vs accuracy, alongside the final console metrics output.")

    # -------------------------------------------------------------
    # STANDALONE GAN SECTION
    # -------------------------------------------------------------
    add_slide(prs, "Model 2: Generative Adversarial Network (GAN) Overview",
              "• Architecture: Generator synthesizes fake images from random noise; Discriminator detects fakes.\n"
              "• Classifier: Shares the same 3-block robust CNN architecture.\n"
              "• Training Loss = Discriminator Loss (BCE Real/Fake) + Generator Adversarial Loss + Classification Loss.\n"
              "• Purpose: Tests if the sharp, realistic edges generated by a GAN improve the classifier's spatial awareness compared to a VAE.")

    add_slide(prs, "GAN Results: Training & Testing Loss",
              "[PLACE IMAGE HERE: GAN Model - Training vs Testing Loss Curve]\n\n"
              "Notes: The dual-line graph tracking the Generator's adversarial and classification loss over 25 epochs.")

    add_slide(prs, "GAN Results: Confusion Matrix & ROC AUC",
              "[PLACE IMAGE HERE: GAN Model - Confusion Matrix Heatmap]\n\n"
              "[PLACE IMAGE HERE: GAN Model - ROC Curve (AUC)]\n\n"
              "Notes: Evaluation of GAN specific classification bounds on the test dataset.")

    add_slide(prs, "GAN Results: Efficiency & Final Metrics",
              "[PLACE IMAGE HERE: GAN Model - Accuracy vs Model Size Line Graph]\n\n"
              "[PLACE IMAGE HERE: GAN Model - Final Metrics Snippet (Accuracy, Dice, Size, Inference)]")

    # -------------------------------------------------------------
    # HYBRID GAN-VAE SECTION
    # -------------------------------------------------------------
    add_slide(prs, "Model 3: Hybrid GAN-VAE Overview",
              "• Architecture: Fuses the Encoder of a VAE with the Generator/Discriminator of a GAN.\n"
              "• Function: VAE Encoder creates a structured latent space, GAN Generator synthesizes highly-detailed images.\n"
              "• Training Loss = MSE + KL Divergence + Adversarial Loss + BCE Classification Loss.\n"
              "• Purpose: Synergize the stable latent features of a VAE with the sharp visual generation of a GAN.")

    add_slide(prs, "Hybrid Results: Training & Testing Loss",
              "[PLACE IMAGE HERE: Hybrid Model - Training vs Testing Loss Curve]\n\n"
              "Notes: The composite loss curve indicating successful synergy and convergence.")

    add_slide(prs, "Hybrid Results: Confusion Matrix & ROC AUC",
              "[PLACE IMAGE HERE: Hybrid Model - Confusion Matrix Heatmap]\n\n"
              "[PLACE IMAGE HERE: Hybrid Model - ROC Curve (AUC)]\n\n"
              "Notes: Shows high true positive rate, low false positive rate, and an exceptional AUC (~0.98).")

    add_slide(prs, "Hybrid Results: Efficiency & Final Metrics",
              "[PLACE IMAGE HERE: Hybrid Model - Accuracy vs Model Size Line Graph]\n\n"
              "[PLACE IMAGE HERE: Hybrid Model - Final Metrics Snippet showing ~93-96% Accuracy]\n\n"
              "Notes: Highlights the optimal balance between a ~4.2MB footprint and clinical-grade accuracy.")

    # -------------------------------------------------------------
    # CONCLUSION
    # -------------------------------------------------------------
    add_slide(prs, "Ablation Study & Edge Deployment",
              "• Quantitative Deployment Analysis:\n"
              "  - Final Parameters: 4.2 MB footprint, ~0.85s inference time per image on a Raspberry Pi 4.\n"
              "  - Synthetic Data Boost: +7% accuracy boost utilizing GAN synthetic generations.\n"
              "• Comparability: Surpasses standard VAE-CNNs (91.5%). While massive GAN-Transformers achieved 96.8%, they are strictly non-deployable on edge devices.\n\n"
              "[PLACE IMAGE HERE: Table 2. Ablation Study Comparison] OR [Table 3. SOTA Comparison]")

    add_slide(prs, "Conclusion & Future Work",
              "• Summary: Successfully developed a 4.2 MB GAN-VAE hybrid capable of 93.2%+ accuracy in sub-second inference on edge devices.\n"
              "• Bridges the gap between high-performance generative medical AI and point-of-care mobile diagnostics.\n"
              "• Limitations: Evaluated on synthetic data to emulate BraTS due to constraints; strictly binary classification scope.\n"
              "• Future Work:\n"
              "  - Validation on real-world global clinical datasets.\n"
              "  - Extending to multi-class tumor grading.\n"
              "  - Integration of federated learning for privacy-preserving, institutional training.")

    prs.save('Comprehensive_Brain_Tumor_Presentation.pptx')
    print("Presentation saved as 'Comprehensive_Brain_Tumor_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
